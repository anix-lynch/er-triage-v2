"""Build pitch.pptx from the slides in pitch.md."""
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

NAVY  = RGBColor(0x0E, 0x2A, 0x47)
GRAY  = RGBColor(0x55, 0x55, 0x55)
LIGHT = RGBColor(0xF4, 0xF6, 0xF8)
RED   = RGBColor(0xC0, 0x39, 0x2B)
GREEN = RGBColor(0x27, 0xAE, 0x60)
AMBER = RGBColor(0xE6, 0x7E, 0x22)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]


def add_textbox(slide, left, top, width, height, text, *, size=18, bold=False, color=NAVY, align=PP_ALIGN.LEFT, font="Helvetica"):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = line
        r.font.name = font
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
    return tb


def add_bullets(slide, left, top, width, height, items, *, size=18, color=NAVY, font="Helvetica"):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(8)
        r = p.add_run()
        r.text = "•  " + item
        r.font.name = font
        r.font.size = Pt(size)
        r.font.color.rgb = color


def header(slide, eyebrow, title):
    add_textbox(slide, 0.6, 0.4, 12, 0.4, eyebrow, size=12, bold=True, color=GRAY)
    add_textbox(slide, 0.6, 0.8, 12, 1.0, title, size=32, bold=True, color=NAVY)


# ── SLIDE 1 — Why we built this ──────────────────────────────────
s = prs.slides.add_slide(BLANK)
header(s, "01  ·  PROBLEM", "Triage today depends on which nurse is at the desk.")
add_bullets(s, 0.6, 2.4, 12, 4, [
    "Same chest pain walks in twice → two different acuity scores, two different waits.",
    "No paper trail of why a tier was assigned.",
    "No flag when a clock is about to be missed (door-to-ECG, last monitored bed).",
    "When something goes wrong, the chart doesn't show what was considered or ruled out.",
])
add_textbox(s, 0.6, 6.2, 12, 0.8,
            "The first 10 minutes of an ED visit are the most consequential — and the least documented.",
            size=18, bold=True, color=RED)


# ── SLIDE 2 — What the prototype shows ───────────────────────────
s = prs.slides.add_slide(BLANK)
header(s, "02  ·  PROTOTYPE", "Side-by-side input and output. Cited. Overridable.")

add_textbox(s, 0.6, 2.0, 6, 0.4, "📥  INPUT (left column)", size=14, bold=True, color=GRAY)
add_bullets(s, 0.6, 2.5, 6, 4, [
    "Vitals, chief complaint, nurse notes",
    "History, allergies, meds, prior visits",
    "Arrival timeline",
    "Add-update form → re-runs assessment",
], size=16)

add_textbox(s, 7.0, 2.0, 6, 0.4, "📤  OUTPUT (right column)", size=14, bold=True, color=GRAY)
add_bullets(s, 7.0, 2.5, 6, 4.5, [
    "Urgency badge + confidence",
    "Constraints: beds, staff, time-target breaches",
    "Hypotheses with confidence bars",
    "Next actions  🔴 Immediate · 🟢 Monitor · 🟡 Escalate",
    "Explanation cites real ESI rule_ids",
    "Follow-up checklist (owner + deadline)",
    "📝 Review note: why this tier, what would flip it",
], size=16)

add_textbox(s, 0.6, 6.7, 12, 0.5,
            "Below the columns: clinician confirms or overrides — every override logged with a reason.",
            size=14, color=GRAY)


out = Path("pitch.pptx")
prs.save(out)
print(f"wrote {out.resolve()}  ({out.stat().st_size:,} bytes, {len(prs.slides)} slides)")
import sys; sys.exit(0)

# ── SLIDE 3 — Proof: asymmetric bias ─────────────────────────────
s = prs.slides.add_slide(BLANK)
header(s, "03  ·  PROOF", "It doesn't blindly trust either signal.")

add_textbox(s, 0.6, 2.0, 12, 0.5,
            "Two adversarial cases. Both correctly classified 🔴 NOW.",
            size=18, color=GRAY)

# Table
rows, cols = 3, 4
left, top, width, height = Inches(0.6), Inches(2.7), Inches(12.1), Inches(2.6)
table = s.shapes.add_table(rows, cols, left, top, width, height).table
table.columns[0].width = Inches(2.2)
table.columns[1].width = Inches(2.5)
table.columns[2].width = Inches(3.5)
table.columns[3].width = Inches(3.9)

headers = ["Case", "Vitals", "Nurse note", "Verdict & why"]
for i, h in enumerate(headers):
    cell = table.cell(0, i)
    cell.text = h
    for p in cell.text_frame.paragraphs:
        for r in p.runs:
            r.font.bold = True
            r.font.size = Pt(13)
            r.font.color.rgb = NAVY

data = [
    ["ER-0052  Walter, 64M", "normal HR/BP/SpO2",
     "“ashen, drowsy, this is not him”",
     "🔴 NOW · medium\nAltered mental status fires ESI-T2 regardless of vitals; amlodipine masks tachycardia"],
    ["ER-0053  Diana, 28F", "shock: HR 142 / SBP 88 / SpO2 91 / T 38.7",
     "“I feel fine, mom is overprotective”",
     "🔴 NOW · high\n‘Appearance is not a vital sign’ — never downgrade shock vitals on self-report"],
]
for r_idx, row in enumerate(data, start=1):
    for c_idx, val in enumerate(row):
        cell = table.cell(r_idx, c_idx)
        cell.text = val
        for p in cell.text_frame.paragraphs:
            for r in p.runs:
                r.font.size = Pt(11)
                r.font.color.rgb = NAVY

add_textbox(s, 0.6, 5.6, 12, 0.5, "By design:", size=14, bold=True, color=GRAY)
add_bullets(s, 0.6, 6.0, 12, 1.4, [
    "Any concerning signal — vitals OR notes OR history — can escalate the tier.",
    "Reassurance never down-triages. Conservative bias breaks ties, surfaced in the review note.",
], size=14)


# ── SLIDE 4 — How it works ───────────────────────────────────────
s = prs.slides.add_slide(BLANK)
header(s, "04  ·  ARCHITECTURE", "How it works.")

flow = (
    "inputs/                          app/                       outputs/\n"
    "patients.json   ──┐\n"
    "ed_state.json   ──┼──►  engine.py (Claude Opus 4.7) ──►  assessments/\n"
    "guidelines.md   ──┘     · prompt-cached system + rules     {case_id}.json\n"
    "                        · tool_use forces 6-panel JSON\n"
    "                        · cites only real rule_ids\n"
    "\n"
    "                        streamlit_app.py  ◄──────────────────┘\n"
    "                        · two-column input / output layout\n"
    "                        · sidebar tier reads from assessment\n"
    "                          (traceable to model decision, not pre-set)"
)
add_textbox(s, 0.6, 2.0, 12.3, 4.0, flow, size=14, font="Menlo", color=NAVY)

add_bullets(s, 0.6, 6.0, 12, 1.4, [
    "12 patient cases  ·  10 baseline + 2 adversarial",
    "18 named ESI v4 rule_ids the model is allowed to cite",
    "~$0.02 per case  ·  8–15 sec latency",
], size=14)


# ── SLIDE 5 — The Ask ────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
header(s, "05  ·  THE ASK", "30-minute call. Pick one ED. Run shadow mode for 4 weeks.")

add_bullets(s, 0.6, 2.4, 12, 4, [
    "Charge nurse + IT lead, 30-minute call.",
    "One ED, four-week observation period.",
    "Shadow mode — assistant generates, nurse decides as today.",
    "We publish the disagreement rate against your senior clinicians' chart reviews.",
    "You decide whether to keep it.",
], size=18)

add_textbox(s, 0.6, 6.0, 12, 1.0,
            "Decision support, not decision.\nCited rules, not opinions.\nLogged overrides, not a black box.",
            size=18, bold=True, color=NAVY)


out = Path("pitch.pptx")
prs.save(out)
print(f"wrote {out.resolve()}  ({out.stat().st_size:,} bytes, {len(prs.slides)} slides)")
